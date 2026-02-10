"""Document processing for multiple file formats."""

import os
from typing import List, Dict, Any, Optional
from pathlib import Path
import re

# Document parsers
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

from utils.config import config
from utils.logger import logger


class DocumentProcessor:
    """Process and chunk documents from various formats."""
    
    def __init__(
        self,
        chunk_size: int = None,
        chunk_overlap: int = None
    ):
        """
        Initialize document processor.
        
        Args:
            chunk_size: Size of text chunks
            chunk_overlap: Overlap between chunks
        """
        self.chunk_size = chunk_size or config.CHUNK_SIZE
        self.chunk_overlap = chunk_overlap or config.CHUNK_OVERLAP
        logger.info(f"Initialized DocumentProcessor (chunk_size={self.chunk_size}, overlap={self.chunk_overlap})")
    
    def process_file(self, file_path: str) -> Dict[str, Any]:
        """
        Process a file and extract text.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Dict with text, chunks, and metadata
        """
        try:
            file_ext = Path(file_path).suffix.lower()
            file_name = Path(file_path).name
            
            logger.info(f"Processing file: {file_name}")
            
            # Extract text based on file type
            if file_ext == '.pdf':
                text = self._extract_pdf(file_path)
            elif file_ext == '.docx':
                text = self._extract_docx(file_path)
            elif file_ext == '.pptx':
                text = self._extract_pptx(file_path)
            elif file_ext == '.xlsx':
                text = self._extract_excel(file_path)
            elif file_ext == '.txt':
                text = self._extract_txt(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")
            
            # Clean text
            text = self._clean_text(text)
            
            # Create chunks
            chunks = self._create_chunks(text)
            
            # Create metadata for each chunk
            metadatas = [
                {
                    'source': file_name,
                    'file_type': file_ext,
                    'chunk_index': i,
                    'total_chunks': len(chunks)
                }
                for i in range(len(chunks))
            ]
            
            logger.info(f"Processed {file_name}: {len(chunks)} chunks created")
            
            return {
                'text': text,
                'chunks': chunks,
                'metadatas': metadatas,
                'file_name': file_name,
                'file_type': file_ext
            }
            
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            raise
    
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF."""
        try:
            text = []
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text.append(f"[Page {page_num + 1}]\n{page_text}")
            
            return "\n\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting PDF: {e}")
            raise
    
    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX."""
        try:
            doc = DocxDocument(file_path)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            logger.error(f"Error extracting DOCX: {e}")
            raise
    
    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX."""
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = [f"[Slide {slide_num + 1}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # More than just the slide number
                    text.append("\n".join(slide_text))
            
            return "\n\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting PPTX: {e}")
            raise
    
    def _extract_excel(self, file_path: str) -> str:
        """Extract text from Excel."""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert dataframe to text representation
                sheet_text = [f"[Sheet: {sheet_name}]"]
                sheet_text.append(f"Columns: {', '.join(df.columns.astype(str))}")
                
                # Add rows (limit to reasonable size)
                for idx, row in df.head(100).iterrows():
                    row_text = " | ".join([f"{col}: {val}" for col, val in row.items()])
                    sheet_text.append(row_text)
                
                # Add summary statistics for numeric columns
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    sheet_text.append("\nSummary Statistics:")
                    for col in numeric_cols:
                        stats = f"{col} - Mean: {df[col].mean():.2f}, Min: {df[col].min():.2f}, Max: {df[col].max():.2f}"
                        sheet_text.append(stats)
                
                text.append("\n".join(sheet_text))
            
            return "\n\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting Excel: {e}")
            raise
    
    def _extract_txt(self, file_path: str) -> str:
        """Extract text from TXT."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error extracting TXT: {e}")
            raise
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        # Remove multiple spaces
        text = re.sub(r'\s+', ' ', text)
        
        # Remove multiple newlines
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        # Strip leading/trailing whitespace
        text = text.strip()
        
        return text
    
    def _create_chunks(self, text: str) -> List[str]:
        """
        Create overlapping chunks from text.
        
        Args:
            text: Input text
            
        Returns:
            List of text chunks
        """
        if not text:
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            # Get chunk
            end = start + self.chunk_size
            chunk = text[start:end]
            
            # Try to break at sentence boundary if not at end
            if end < text_length:
                # Look for sentence endings
                last_period = chunk.rfind('. ')
                last_newline = chunk.rfind('\n')
                last_break = max(last_period, last_newline)
                
                if last_break > self.chunk_size * 0.5:  # Only break if we're past halfway
                    chunk = chunk[:last_break + 1]
                    end = start + last_break + 1
            
            chunks.append(chunk.strip())
            
            # Move start position with overlap
            start = end - self.chunk_overlap
            
            # Prevent infinite loop
            if end >= text_length:
                break
        
        # Remove empty chunks
        chunks = [c for c in chunks if c.strip()]
        
        return chunks


# Create singleton instance
document_processor = DocumentProcessor()
